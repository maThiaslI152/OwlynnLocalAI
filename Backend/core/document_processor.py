import os
from typing import Dict, Any, Optional, List, Tuple
import fitz  # PyMuPDF
import docx
from striprtf.striprtf import rtf_to_text
import pandas as pd
from pptx import Presentation
import json
import yaml
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image
import cv2
import numpy as np
from transformers import BlipProcessor, BlipForConditionalGeneration
import spacy
from core.config import settings

class DocumentProcessor:
    def __init__(self):
        # Initialize BLIP for image captioning
        self.blip_processor = BlipProcessor.from_pretrained(settings.BLIP_MODEL)
        self.blip_model = BlipForConditionalGeneration.from_pretrained(settings.BLIP_MODEL)
        
        # Initialize spaCy for text chunking
        self.nlp = spacy.load("en_core_web_sm")
        
        # Create necessary directories
        os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
        os.makedirs(settings.CACHE_DIR, exist_ok=True)
    
    def process_file(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process a file and return its content and metadata"""
        file_ext = os.path.splitext(file_path)[1].lower()
        content = ""
        metadata = {
            "filename": os.path.basename(file_path),
            "file_type": file_ext,
            "file_size": os.path.getsize(file_path)
        }
        
        try:
            if file_ext in settings.SUPPORTED_EXTENSIONS["text"]:
                content = self._process_text_file(file_path)
            elif file_ext in settings.SUPPORTED_EXTENSIONS["documents"]:
                content = self._process_document(file_path, file_ext)
            elif file_ext in settings.SUPPORTED_EXTENSIONS["spreadsheets"]:
                content = self._process_spreadsheet(file_path, file_ext)
            elif file_ext in settings.SUPPORTED_EXTENSIONS["presentations"]:
                content = self._process_presentation(file_path)
            elif file_ext in settings.SUPPORTED_EXTENSIONS["code"]:
                content = self._process_code_file(file_path, file_ext)
            elif file_ext in settings.SUPPORTED_EXTENSIONS["images"]:
                content, image_metadata = self._process_image(file_path)
                metadata.update(image_metadata)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            # Chunk the content
            chunks = self._chunk_text(content)
            metadata["chunks"] = chunks
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Error processing file {file_path}: {str(e)}")
    
    def _process_text_file(self, file_path: str) -> str:
        """Process plain text files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _process_document(self, file_path: str, file_ext: str) -> str:
        """Process PDF, DOCX, and RTF files"""
        if file_ext == '.pdf':
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        elif file_ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_ext == '.rtf':
            with open(file_path, 'r', encoding='utf-8') as f:
                return rtf_to_text(f.read())
    
    def _process_spreadsheet(self, file_path: str, file_ext: str) -> str:
        """Process CSV and Excel files"""
        if file_ext == '.csv':
            df = pd.read_csv(file_path)
        else:  # .xlsx
            df = pd.read_excel(file_path)
        return df.to_string()
    
    def _process_presentation(self, file_path: str) -> str:
        """Process PowerPoint files"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    
    def _process_code_file(self, file_path: str, file_ext: str) -> str:
        """Process code files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
        if file_ext in ['.json']:
            # Pretty print JSON
            return json.dumps(json.loads(content), indent=2)
        elif file_ext in ['.yml', '.yaml']:
            # Pretty print YAML
            return yaml.dump(yaml.safe_load(content), default_flow_style=False)
        elif file_ext in ['.html', '.xml']:
            # Clean HTML/XML
            soup = BeautifulSoup(content, 'lxml')
            return soup.prettify()
        else:
            return content
    
    def _process_image(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process image files with OCR and captioning"""
        # Read image
        image = Image.open(file_path)
        
        # Perform OCR
        ocr_text = pytesseract.image_to_string(image, lang='+'.join(settings.TESSERACT_LANGUAGES))
        
        # Generate image caption
        inputs = self.blip_processor(image, return_tensors="pt")
        out = self.blip_model.generate(**inputs)
        caption = self.blip_processor.decode(out[0], skip_special_tokens=True)
        
        # Extract image metadata
        metadata = {
            "dimensions": image.size,
            "format": image.format,
            "mode": image.mode,
            "caption": caption
        }
        
        return f"OCR Text:\n{ocr_text}\n\nCaption:\n{caption}", metadata
    
    def _chunk_text(self, text: str, chunk_size: int = 1000) -> List[str]:
        """Split text into chunks using spaCy"""
        doc = self.nlp(text)
        chunks = []
        current_chunk = []
        current_size = 0
        
        for sent in doc.sents:
            sent_text = sent.text.strip()
            sent_size = len(sent_text)
            
            if current_size + sent_size > chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = [sent_text]
                current_size = sent_size
            else:
                current_chunk.append(sent_text)
                current_size += sent_size
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks 